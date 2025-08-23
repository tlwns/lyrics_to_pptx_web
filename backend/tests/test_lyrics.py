"""
Test suite for the lyrics module.
"""

from io import BytesIO
from typing import cast

from fastapi import status
from pptx import Presentation
from pptx.shapes.autoshape import Shape
import pytest

from tests.data import a_to_z_lyrics_sample


def test_generate_empty_lyrics(client):
    """Test generating a PowerPoint with empty lyrics."""
    response = client.post(
        "/generate",
        json={"lyrics": "", "filename": "test_lyrics", "background_option": "NONE"},
    )

    assert response.status_code == status.HTTP_400_BAD_REQUEST
    assert response.json() == {"detail": "Lyrics cannot be empty."}


def test_generate_empty_filename(client):
    """Test generating a PowerPoint with empty filename."""
    response = client.post(
        "/generate",
        json={"lyrics": "Sample lyrics", "filename": "", "background_option": "NONE"},
    )

    assert response.status_code == status.HTTP_200_OK
    assert response.headers["Content-Disposition"] == "attachment; filename=lyrics.pptx"
    assert response.headers["Content-Type"] == (
        "application/vnd.openxmlformats-officedocument." "presentationml.presentation"
    )


def test_generate_invalid_background_option(client):
    """Test generating a PowerPoint with an invalid background option."""
    response = client.post(
        "/generate",
        json={
            "lyrics": "Sample lyrics",
            "filename": "test_lyrics",
            "background_option": "INVALID_OPTION",
        },
    )

    assert response.status_code == status.HTTP_422_UNPROCESSABLE_ENTITY


def test_generate_valid_lyrics(client):
    """Test generating a PowerPoint with valid lyrics."""
    response = client.post(
        "/generate",
        json={
            "lyrics": "Sample lyrics",
            "filename": "test_lyrics",
            "background_option": "NONE",
        },
    )

    assert response.status_code == status.HTTP_200_OK
    assert response.headers["Content-Disposition"] == "attachment; filename=test_lyrics"
    assert response.headers["Content-Type"] == (
        "application/vnd.openxmlformats-officedocument." "presentationml.presentation"
    )

    content = response.content
    pres = Presentation(BytesIO(content))
    assert len(pres.slides) == 2  # Lyrics and empty final slide
    shape: Shape = cast(Shape, pres.slides[0].shapes[0])
    assert shape.text == "\nSample lyrics"


@pytest.mark.parametrize(
    "input_data, expected",
    [(a_to_z_lyrics_sample.data["input"], a_to_z_lyrics_sample.data["expected"])],
)
def test_generate_genius_lyrics_valid(client, input_data, expected):
    """Test generating a PowerPoint with Genius lyrics."""
    response = client.post(
        "/generate",
        json={
            "lyrics": input_data["lyrics"],
            "filename": input_data["filename"],
            "background_option": input_data["background_option"],
        },
    )

    assert response.status_code == status.HTTP_200_OK
    assert (
        response.headers["Content-Disposition"]
        == f"attachment; filename={input_data['filename']}"
    )
    assert response.headers["Content-Type"] == (
        "application/vnd.openxmlformats-officedocument." "presentationml.presentation"
    )

    content = response.content
    pres = Presentation(BytesIO(content))
    assert len(pres.slides) == expected["num_slides"]

    for i, slide in enumerate(pres.slides):
        # Check content of each slide

        assert len(slide.shapes) == 1
        shape: Shape = cast(Shape, slide.shapes[0])
        if expected["slides_text"][i]:
            # Check if the slide text matches the expected text
            assert shape.text == "\n" + "\n".join(expected["slides_text"][i])
        else:
            # If the expected text is empty, the slide text should also be empty
            assert shape.text == ""
