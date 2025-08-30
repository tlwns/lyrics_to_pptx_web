"""
Service layer for processing lyrics and generating PowerPoint presentations.
"""

from io import BytesIO

from pptx import Presentation
from pptx.util import Length
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from lyrics.utils import is_hangul, split_lyrics, split_lines
from lyrics.schemas import BackgroundOption

DEFAULT_BACKGROUND_IMAGE = "app/static/gift_background.jpg"


def generate_pptx_in_memory(lyrics: str, background_option: BackgroundOption) -> bytes:
    """
    Generate a PowerPoint presentation in memory with the given lyrics.

    Args:
        lyrics (str): The lyrics to include in the PowerPoint slides.
        background_option (BackgroundOption): The background option for the slides.

    Returns:
        bytes: The bytes object of the generated PowerPoint file.
    """

    if not lyrics.strip():
        raise ValueError("Lyrics cannot be empty.")

    buffer = BytesIO()
    pres = build_pptx(lyrics, background_option)
    pres.save(buffer)
    buffer.seek(0)  # Reset the buffer position to the beginning
    return buffer.read1()


def build_pptx(lyrics: str, background_option: BackgroundOption) -> Presentation:  # type: ignore
    """
    Build a PowerPoint presentation with the given lyrics.

    Args:
        lyrics (str): The lyrics to include in the PowerPoint slides.
        background_option (BackgroundOption): The background option for the slides.

    Returns:
        Presentation: A Presentation object containing the slides.
    """

    pres = Presentation(
        # import from static files or a default template
        "app/static/blank-16x9.pptx"
    )

    lyrics_sections = split_lyrics(lyrics)

    for section in lyrics_sections:
        if section == "////":
            # Add a black slide if new song is indicated by "////"
            slide = pres.slides.add_slide(pres.slide_layouts[5])
            if background_option is BackgroundOption.GIFT:
                slide.shapes.add_picture(
                    DEFAULT_BACKGROUND_IMAGE,
                    Inches(0),
                    Inches(0),
                    width=pres.slide_width,
                    height=pres.slide_height,
                )
            continue

        lines = section.splitlines()
        lines = [line.strip() for line in lines]

        # If there is more than 5 lines in a single section, split them into two slides
        parts = split_lines(lines)

        for part in parts:
            # Create a new slide with a blank layout
            # and add the background image if provided
            slide = pres.slides.add_slide(pres.slide_layouts[6])
            if background_option is BackgroundOption.GIFT:
                slide.shapes.add_picture(
                    DEFAULT_BACKGROUND_IMAGE,
                    Inches(0),
                    Inches(0),
                    width=pres.slide_width,
                    height=pres.slide_height,
                )

            # Add a textbox with the lyrics
            box = slide.shapes.add_textbox(
                Inches(0),
                Inches(2),
                Length(pres.slide_width if pres.slide_width else 0),
                Length(pres.slide_height if pres.slide_height else 0),
            )
            frame = box.text_frame

            for line in part:
                para = frame.add_paragraph()
                para.font.size = Pt(48)
                para.line_spacing = 2
                para.text = line
                para.alignment = PP_ALIGN.CENTER
                if is_hangul(line):
                    para.font.name = "Spoqa Han Sans Neo"
                else:
                    para.font.name = "Lexend"

    # Add a final blank slide
    slide = pres.slides.add_slide(pres.slide_layouts[5])
    if background_option is BackgroundOption.GIFT:
        slide.shapes.add_picture(
            DEFAULT_BACKGROUND_IMAGE,
            Inches(0),
            Inches(0),
            width=pres.slide_width,
            height=pres.slide_height,
        )

    return pres
